#!/usr/bin/env python3
from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from matplotlib.lines import Line2D
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_MD = ROOT / "docs" / "presentations" / "presentation.md"
OUTPUT_PPTX = ROOT / "docs" / "presentations" / "final_presentation.pptx"
VRAM_CHART = ROOT / "docs" / "figures" / "vram_compression.png"
RUNTIME_CHART = ROOT / "docs" / "figures" / "runtime_quality.png"
DRIFT_CHART = ROOT / "docs" / "figures" / "temporal_drift.png"

DATASET_CANDIDATES = [
    ROOT / "results" / "combined" / "combined_comparison_dataset.csv",
    ROOT / "combined_comparison_dataset.csv",
]

METHOD_ORDER_BASE = [
    "BF16",
    "RTN_INT4",
    "FLOWCACHE_SOFT_PRUNE_INT4",
    "SPATIAL_MIXED_FG_QUAROT_KV_INT4_BG_RTN_INT2",
]

COLOR_MAP = {
    "BF16": "#111111",
    "RTN_INT4": "#2F6BDE",
    "PRQ_INT4": "#D1495B",
    "PRQ_INT2": "#D1495B",
    "FLOWCACHE_SOFT_PRUNE_INT4": "#228B22",
    "SPATIAL_MIXED_FG_QUAROT_KV_INT4_BG_RTN_INT2": "#7A7A7A",
}

LABEL_MAP = {
    "BF16": "BF16",
    "RTN_INT4": "RTN INT4",
    "PRQ_INT4": "PRQ INT4",
    "PRQ_INT2": "PRQ INT2",
    "FLOWCACHE_SOFT_PRUNE_INT4": "Custom FlowCache\nSoft-Prune INT4",
    "SPATIAL_MIXED_FG_QUAROT_KV_INT4_BG_RTN_INT2": "Spatial Mixed\nFailure",
}

ACCENT_GREEN = RGBColor(34, 139, 34)
TEXT_DARK = RGBColor(27, 38, 59)
TEXT_MUTED = RGBColor(68, 84, 106)
LIGHT_BORDER = RGBColor(196, 205, 213)


@dataclass
class SlideSpec:
    number: int
    title: str
    bullet_points: list[str]
    visuals_required: str
    speaker_notes: str
    dashboard_integration: str


def find_dataset_path() -> Path:
    for candidate in DATASET_CANDIDATES:
        if candidate.exists():
            return candidate
    candidates = "\n".join(str(path) for path in DATASET_CANDIDATES)
    raise FileNotFoundError(f"Could not locate combined comparison dataset. Tried:\n{candidates}")


def first_valid(series: pd.Series) -> float | None:
    values = series.dropna()
    if values.empty:
        return None
    return float(values.iloc[0])


def to_numeric(df: pd.DataFrame, columns: Iterable[str]) -> pd.DataFrame:
    for column in columns:
        if column in df.columns:
            df[column] = pd.to_numeric(df[column], errors="coerce")
    return df


def resolve_methods(df: pd.DataFrame) -> list[str]:
    available = set(df["method"].dropna().astype(str))
    prq_method = "PRQ_INT4" if "PRQ_INT4" in available else "PRQ_INT2"
    methods = ["BF16", "RTN_INT4", prq_method, *METHOD_ORDER_BASE[2:]]
    missing = [method for method in methods if method not in available]
    if missing:
        raise ValueError(f"Required narrative methods are missing from the dataset: {missing}")
    return methods


def summarize_benchmark(df: pd.DataFrame, benchmark: str, methods: list[str]) -> pd.DataFrame:
    subset = df[df["benchmark"] == benchmark].copy()
    rows: list[dict[str, float | str]] = []
    imaging_col = f"{benchmark}_imaging_quality_agg"
    fallback_imaging_col = f"{benchmark}_imaging_quality"
    drift_col = f"{benchmark}_drift_last_imaging_quality"

    for method in methods:
        method_rows = subset[subset["method"] == method]
        if method_rows.empty:
            continue

        peak_bytes = method_rows["peak_vram_bytes"].dropna()
        peak_mb = method_rows["peak_vram_mb"].dropna()
        runtime = first_valid(method_rows["avg_runtime_s_per_prompt"])
        compression = first_valid(method_rows["compression_ratio"])
        imaging = first_valid(method_rows[imaging_col]) if imaging_col in method_rows.columns else None
        if imaging is None and fallback_imaging_col in method_rows.columns:
            imaging = first_valid(method_rows[fallback_imaging_col])
            if imaging is not None and imaging > 1.0:
                imaging /= 100.0
        drift = first_valid(method_rows[drift_col]) if drift_col in method_rows.columns else None

        if not peak_bytes.empty:
            peak_vram_gb = float(peak_bytes.max()) / (1024**3)
        elif not peak_mb.empty:
            peak_vram_gb = float(peak_mb.max()) / 1024.0
        else:
            peak_vram_gb = None

        rows.append(
            {
                "benchmark": benchmark,
                "method": method,
                "method_label": LABEL_MAP.get(method, method),
                "color": COLOR_MAP.get(method, "#4C78A8"),
                "compression_ratio": compression,
                "peak_vram_gb": peak_vram_gb,
                "runtime_s": runtime,
                "imaging_quality": imaging,
                "drift_last": drift,
            }
        )

    summary = pd.DataFrame(rows)
    if summary.empty:
        raise ValueError(f"No summary rows produced for benchmark={benchmark}")
    return summary


def load_dataset() -> tuple[pd.DataFrame, dict[str, pd.DataFrame], list[str], Path]:
    dataset_path = find_dataset_path()
    df = pd.read_csv(dataset_path)
    numeric_columns = [
        "peak_vram_bytes",
        "peak_vram_mb",
        "compression_ratio",
        "avg_runtime_s_per_prompt",
        "moviegen_imaging_quality",
        "moviegen_imaging_quality_agg",
        "moviegen_drift_last_imaging_quality",
        "storyeval_imaging_quality",
        "storyeval_imaging_quality_agg",
        "storyeval_drift_last_imaging_quality",
    ]
    df = to_numeric(df, numeric_columns)
    methods = resolve_methods(df)
    filtered = df[df["method"].isin(methods)].copy()
    summaries = {
        "moviegen": summarize_benchmark(filtered, "moviegen", methods),
        "storyeval": summarize_benchmark(filtered, "storyeval", methods),
    }
    return filtered, summaries, methods, dataset_path


def parse_presentation_markdown(markdown_path: Path) -> list[SlideSpec]:
    if not markdown_path.exists():
        raise FileNotFoundError(f"presentation markdown not found: {markdown_path}")

    text = markdown_path.read_text(encoding="utf-8")
    slide_numbers = [int(value) for value in re.findall(r"^## Slide (\d+)\s*$", text, flags=re.MULTILINE)]
    blocks = re.split(r"^## Slide \d+\s*$", text, flags=re.MULTILINE)[1:]

    if len(slide_numbers) != len(blocks):
        raise ValueError("Failed to align slide headers with slide bodies while parsing presentation.md")

    slides: list[SlideSpec] = []
    for number, block in zip(slide_numbers, blocks):
        title = ""
        bullets: list[str] = []
        visuals_required = ""
        speaker_notes = ""
        dashboard_integration = ""
        mode: str | None = None

        for raw_line in block.strip().splitlines():
            line = raw_line.rstrip()
            stripped = line.strip()
            if not stripped:
                continue

            title_match = re.match(r"^- \*\*Slide Title:\*\*\s*(.*)$", stripped)
            if title_match:
                title = title_match.group(1).strip()
                mode = None
                continue
            content_match = re.match(r"^- \*\*Slide Content \(Bullet Points\):\*\*\s*(.*)$", stripped)
            if content_match:
                mode = "bullets"
                remainder = content_match.group(1).strip()
                if remainder:
                    bullets.append(remainder)
                continue
            visuals_match = re.match(r"^- \*\*Visuals Required:\*\*\s*(.*)$", stripped)
            if visuals_match:
                visuals_required = visuals_match.group(1).strip()
                mode = None
                continue
            notes_match = re.match(r"^- \*\*Speaker Notes \(Script\):\*\*\s*(.*)$", stripped)
            if notes_match:
                speaker_notes = notes_match.group(1).strip().strip("“”\"")
                mode = None
                continue
            dashboard_match = re.match(r"^- \*\*Dashboard Integration:\*\*\s*(.*)$", stripped)
            if dashboard_match:
                dashboard_integration = dashboard_match.group(1).strip()
                mode = None
                continue
            if mode == "bullets" and stripped.startswith("- "):
                bullets.append(stripped[2:].strip())

        if not title:
            raise ValueError(f"Slide {number} has no title in presentation.md")

        slides.append(
            SlideSpec(
                number=number,
                title=title,
                bullet_points=bullets,
                visuals_required=visuals_required,
                speaker_notes=speaker_notes,
                dashboard_integration=dashboard_integration,
            )
        )

    return slides


def apply_plot_theme() -> None:
    sns.set_theme(style="whitegrid", context="talk")
    plt.rcParams["figure.dpi"] = 300
    plt.rcParams["savefig.dpi"] = 300
    plt.rcParams["font.family"] = "DejaVu Sans"


def add_value_labels(ax, bars, formatter: str, dy: float) -> None:
    for bar in bars:
        height = bar.get_height()
        if pd.isna(height):
            continue
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            height + dy,
            format(height, formatter),
            ha="center",
            va="bottom",
            fontsize=10,
            color="#222222",
        )


def generate_vram_compression_chart(moviegen: pd.DataFrame) -> None:
    apply_plot_theme()
    ordered = moviegen.copy()
    methods = ordered["method"].tolist()
    labels = [LABEL_MAP.get(method, method) for method in methods]
    colors = [COLOR_MAP.get(method, "#4C78A8") for method in methods]
    x = list(range(len(methods)))
    width = 0.36

    fig, ax1 = plt.subplots(figsize=(12, 6.8))
    ax2 = ax1.twinx()

    vram_bars = ax1.bar(
        [value - width / 2 for value in x],
        ordered["peak_vram_gb"],
        width=width,
        color=colors,
        alpha=0.92,
        edgecolor="white",
        linewidth=1.2,
    )
    compression_bars = ax2.bar(
        [value + width / 2 for value in x],
        ordered["compression_ratio"],
        width=width,
        color=colors,
        alpha=0.35,
        edgecolor=colors,
        linewidth=1.4,
        hatch="//",
    )

    flowcache_idx = methods.index("FLOWCACHE_SOFT_PRUNE_INT4")
    for bar_group in (vram_bars, compression_bars):
        bar_group[flowcache_idx].set_edgecolor("#C08A00")
        bar_group[flowcache_idx].set_linewidth(2.5)

    ax1.set_title("MovieGen Systems Boundary Points: Peak VRAM vs. KV Compression", pad=16, fontsize=20)
    ax1.set_ylabel("Peak VRAM (GB)")
    ax2.set_ylabel("Compression Ratio (x)")
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=11)
    ax1.set_ylim(0, max(ordered["peak_vram_gb"]) * 1.25)
    ax2.set_ylim(0, max(ordered["compression_ratio"]) * 1.35)

    add_value_labels(ax1, vram_bars, ".2f", dy=0.18)
    add_value_labels(ax2, compression_bars, ".2f", dy=0.08)

    metric_legend = [
        Patch(facecolor="#6D8AC5", alpha=0.92, label="Peak VRAM (GB)"),
        Patch(facecolor="#6D8AC5", alpha=0.35, edgecolor="#6D8AC5", hatch="//", label="Compression Ratio (x)"),
        Line2D([0], [0], color="#C08A00", linewidth=2.5, label="Highlighted: Our custom FlowCache adaptation"),
    ]
    ax1.legend(handles=metric_legend, loc="upper right", frameon=True)
    ax1.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)
    fig.tight_layout()
    fig.savefig(VRAM_CHART, bbox_inches="tight")
    plt.close(fig)


def generate_runtime_quality_chart(moviegen: pd.DataFrame) -> None:
    apply_plot_theme()
    fig, ax = plt.subplots(figsize=(11.5, 6.5))

    for _, row in moviegen.iterrows():
        marker = "o"
        size = 220
        edge = "white"
        linewidth = 1.6
        if row["method"] == "FLOWCACHE_SOFT_PRUNE_INT4":
            marker = "D"
            size = 280
            edge = "#C08A00"
            linewidth = 2.0

        ax.scatter(
            row["runtime_s"],
            row["imaging_quality"],
            s=size,
            c=row["color"],
            marker=marker,
            edgecolors=edge,
            linewidths=linewidth,
            zorder=3,
        )
        ax.annotate(
            LABEL_MAP.get(row["method"], row["method"]).replace("\n", " "),
            xy=(row["runtime_s"], row["imaging_quality"]),
            xytext=(8, 8),
            textcoords="offset points",
            fontsize=10,
            color="#222222",
        )

    ax.set_title("MovieGen Systems vs. Quality Trade-off", pad=16, fontsize=20)
    ax.set_xlabel("Runtime per Prompt (s)")
    ax.set_ylabel("Imaging Quality")
    ax.grid(True, linestyle="--", alpha=0.35)
    ax.set_xlim(left=0)
    ax.set_ylim(0.35, 0.78)

    legend_handles = [
        Line2D([0], [0], marker="o", color="w", label="BF16", markerfacecolor=COLOR_MAP["BF16"], markersize=12),
        Line2D([0], [0], marker="o", color="w", label="RTN INT4", markerfacecolor=COLOR_MAP["RTN_INT4"], markersize=12),
        Line2D(
            [0],
            [0],
            marker="o",
            color="w",
            label="PRQ",
            markerfacecolor=COLOR_MAP.get("PRQ_INT4", COLOR_MAP["PRQ_INT2"]),
            markersize=12,
        ),
        Line2D(
            [0],
            [0],
            marker="D",
            color="w",
            label="Our custom FlowCache adaptation",
            markerfacecolor=COLOR_MAP["FLOWCACHE_SOFT_PRUNE_INT4"],
            markeredgecolor="#C08A00",
            markeredgewidth=1.4,
            markersize=12,
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            color="w",
            label="Spatial mixed failure",
            markerfacecolor=COLOR_MAP["SPATIAL_MIXED_FG_QUAROT_KV_INT4_BG_RTN_INT2"],
            markersize=12,
        ),
    ]
    ax.legend(handles=legend_handles, loc="lower right", frameon=True)
    fig.tight_layout()
    fig.savefig(RUNTIME_CHART, bbox_inches="tight")
    plt.close(fig)


def generate_temporal_drift_chart(storyeval: pd.DataFrame) -> None:
    apply_plot_theme()
    fig, ax = plt.subplots(figsize=(11.5, 6.3))
    methods = storyeval["method"].tolist()
    labels = [LABEL_MAP.get(method, method) for method in methods]
    colors = [COLOR_MAP.get(method, "#4C78A8") for method in methods]

    bars = ax.bar(labels, storyeval["drift_last"], color=colors, edgecolor="white", linewidth=1.2)
    flowcache_idx = methods.index("FLOWCACHE_SOFT_PRUNE_INT4")
    bars[flowcache_idx].set_edgecolor("#C08A00")
    bars[flowcache_idx].set_linewidth(2.5)

    ax.set_title("StoryEval Terminal Temporal Drift", pad=16, fontsize=20)
    ax.set_ylabel("Drift Last Imaging Quality")
    ax.set_ylim(0.35, 0.75)
    ax.grid(True, axis="y", linestyle="--", alpha=0.3)

    for bar in bars:
        height = bar.get_height()
        ax.text(
            bar.get_x() + bar.get_width() / 2.0,
            height + 0.008,
            f"{height:.3f}",
            ha="center",
            va="bottom",
            fontsize=10,
        )

    fig.tight_layout()
    fig.savefig(DRIFT_CHART, bbox_inches="tight")
    plt.close(fig)


def generate_visuals(summaries: dict[str, pd.DataFrame]) -> None:
    generate_vram_compression_chart(summaries["moviegen"])
    generate_runtime_quality_chart(summaries["moviegen"])
    generate_temporal_drift_chart(summaries["storyeval"])


def add_title_box(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.7))
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(25)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TEXT_DARK
    paragraph.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=Inches(0.55),
        top=Inches(1.0),
        width=Inches(12.15),
        height=Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BORDER
    line.line.color.rgb = LIGHT_BORDER


def add_body_bullets(slide, bullets: list[str], has_visuals: bool) -> None:
    left = Inches(0.7)
    top = Inches(1.25)
    width = Inches(5.1 if has_visuals else 12.0)
    height = Inches(5.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(19 if has_visuals else 21)
        paragraph.font.color.rgb = TEXT_DARK
        paragraph.space_after = Pt(10)
        paragraph.space_before = Pt(2)


def select_visuals(slide_spec: SlideSpec) -> list[Path]:
    title = slide_spec.title.lower()
    visuals = slide_spec.visuals_required.lower()
    assets: list[Path] = []

    if "systems boundary" in title or "vram" in visuals or "compression" in visuals:
        assets.append(VRAM_CHART)
    if "quality and drift" in title:
        assets.extend([RUNTIME_CHART, DRIFT_CHART])
    elif "runtime" in visuals or "scatter" in visuals or "systems vs. quality" in visuals:
        assets.append(RUNTIME_CHART)
    elif "drift" in title or "temporal" in title or "drift" in visuals:
        assets.append(DRIFT_CHART)

    deduped: list[Path] = []
    for asset in assets:
        if asset.exists() and asset not in deduped:
            deduped.append(asset)
    return deduped


def add_visuals(slide, image_paths: list[Path]) -> None:
    if not image_paths:
        return

    if len(image_paths) == 1:
        slide.shapes.add_picture(str(image_paths[0]), Inches(6.0), Inches(1.35), width=Inches(6.7))
        return

    picture_width = Inches(6.2)
    top_positions = [Inches(1.3), Inches(4.1)]
    for image_path, top in zip(image_paths[:2], top_positions):
        slide.shapes.add_picture(str(image_path), Inches(6.15), top, width=picture_width)


def add_slide_notes(slide, speaker_notes: str, dashboard_integration: str) -> None:
    notes_parts = []
    speaker_notes = speaker_notes.strip()
    dashboard_integration = dashboard_integration.strip()
    if speaker_notes:
        notes_parts.append(f"Speaker Notes:\n{speaker_notes}")
    if dashboard_integration and dashboard_integration.lower() != "no dashboard switch on this slide.":
        notes_parts.append(f"Dashboard Integration:\n{dashboard_integration}")
    elif dashboard_integration:
        notes_parts.append(f"Dashboard Integration:\n{dashboard_integration}")

    notes_text = "\n\n".join(notes_parts).strip()
    notes_slide = slide.notes_slide
    for placeholder in notes_slide.placeholders:
        if getattr(placeholder, "name", "").lower().startswith("notes"):
            placeholder.text = notes_text
            return
    raise RuntimeError("Unable to locate the speaker-notes placeholder on the notes slide.")


def build_title_slide(prs: Presentation, slide_spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = slide_spec.title
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TEXT_DARK

    subtitle_shape.text = "\n".join(f"• {bullet}" for bullet in slide_spec.bullet_points)
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = TEXT_MUTED

    add_slide_notes(slide, slide_spec.speaker_notes, slide_spec.dashboard_integration)


def build_content_slide(prs: Presentation, slide_spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, slide_spec.title)
    visuals = select_visuals(slide_spec)
    add_body_bullets(slide, slide_spec.bullet_points, has_visuals=bool(visuals))
    add_visuals(slide, visuals)
    add_slide_notes(slide, slide_spec.speaker_notes, slide_spec.dashboard_integration)


def build_presentation(slides: list[SlideSpec]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_spec in slides:
        if slide_spec.number == 1:
            build_title_slide(prs, slide_spec)
        else:
            build_content_slide(prs, slide_spec)

    prs.save(OUTPUT_PPTX)


def main() -> None:
    _, summaries, methods, dataset_path = load_dataset()
    slides = parse_presentation_markdown(PRESENTATION_MD)
    generate_visuals(summaries)
    build_presentation(slides)

    print(f"dataset_path={dataset_path}")
    print(f"selected_methods={methods}")
    print(f"presentation_md={PRESENTATION_MD}")
    print(f"generated_plot={VRAM_CHART}")
    print(f"generated_plot={RUNTIME_CHART}")
    print(f"generated_plot={DRIFT_CHART}")
    print(f"generated_pptx={OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
